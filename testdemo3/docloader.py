import os
import textract
import fitz
import docx
import pandas as pd
from pptx import Presentation
from PIL import Image
import pytesseract

class DocumentLoader:
    def load_documents(self):
        # Load documents from CompanyData folder
        for filename in os.listdir("CompanyData"):
            file_path = os.path.join("CompanyData", filename)
            text = self.extract_text_from_file(file_path)
            self.documents.append(text)

    def extract_text_from_file(self, file_path: str) -> str:
        if file_path.endswith(".txt"):
            with open(file_path, "r") as f:
                return f.read()
        elif file_path.endswith(".pdf"):
            with fitz.open(file_path) as pdf:
                text = ""
                for page in pdf:
                    text += page.get_text()
                return text
        elif file_path.endswith(".docx"):
            doc = docx.Document(file_path)
            text = ""
            for para in doc.paragraphs:
                text += para.text
            return text
        elif file_path.endswith('.csv'):
            df = pd.read_csv(file_path)
            return df.to_string()
        elif file_path.endswith('.pptx'):
            return self.extract_text_from_pptx(file_path)
        elif file_path.endswith('.jpg') or file_path.endswith('.png'):
            return self.extract_text_from_image(file_path)
        else:
            text = textract.process(file_path)
            return text.decode("utf-8")

    def extract_text_from_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def extract_text_from_image(self, file_path: str) -> str:
        image = Image.open(file_path)
        return pytesseract.image_to_string(image)
